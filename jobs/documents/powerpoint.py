"""jobs/documents/powerpoint.py — Read and create PowerPoint .pptx files."""
import logging
import re

log = logging.getLogger(__name__)


def read_pptx(path: str) -> str:
    from pptx import Presentation
    try:
        prs = Presentation(path)
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"Slide {i}:")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(f"  {text}")
        return "\n".join(lines)
    except Exception as exc:
        log.error("read_pptx failed: %s", exc)
        return f"Error reading PowerPoint file: {exc}"


def create_pptx(slides: list, path: str) -> bool:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    try:
        prs = Presentation()
        layout = prs.slide_layouts[1]  # title + content
        for slide_data in slides:
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = slide_data.get("title", "")
            body = slide.placeholders[1]
            body.text = slide_data.get("content", "")
        prs.save(path)
        return True
    except Exception as exc:
        log.error("create_pptx failed: %s", exc)
        return False


def run(message: str = None) -> str:
    if not message:
        return "PowerPoint skill ready."
    match = re.search(r'[\w/~.-]+\.pptx?', message, re.IGNORECASE)
    if not match:
        return "PowerPoint skill ready. Provide a file path to read a presentation."
    path = match.group(0).replace("~", __import__("os").path.expanduser("~"))
    text = read_pptx(path)
    if not text:
        return f"No text extracted from {path}."
    preview = text[:1000]
    suffix = f"\n\n[{len(text)} chars total]" if len(text) > 1000 else ""
    return f"PowerPoint: {path}\n\n{preview}{suffix}"
