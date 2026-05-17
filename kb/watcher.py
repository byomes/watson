"""
watcher.py — Knowledge Base file watcher for Windows desktop.
Watches F:\Knowledge_Database\_inbox for new files.
Converts to .txt/.md, saves to F:\Knowledge_Database\_library,
ingests into OpenWebUI, deletes original from inbox.

Run on startup via Task Scheduler.
"""

import os
import sys
import time
import shutil
import logging
import requests
from pathlib import Path
from dotenv import load_dotenv

load_dotenv(r"D:\OneDrive\Claude\agents\watson\.env")

INBOX = Path(r"F:\Knowledge_Database\_inbox")
LIBRARY = Path(r"F:\Knowledge_Database\_library")
OPENWEBUI_URL = os.getenv("OPENWEBUI_URL", "http://localhost:3000")
OPENWEBUI_API_KEY = os.getenv("OPENWEBUI_API_KEY")
KNOWLEDGE_COLLECTION = "Personal Library"
CACHE_FILE = Path(__file__).parent / ".collection_id_cache.json"

_collection_id = None


def _load_cached_collection_id() -> str | None:
    try:
        if CACHE_FILE.exists():
            import json
            data = json.loads(CACHE_FILE.read_text(encoding="utf-8"))
            return data.get("collection_id")
    except Exception:
        pass
    return None


def _save_cached_collection_id(collection_id: str):
    import json
    CACHE_FILE.write_text(
        json.dumps({"collection_id": collection_id, "name": KNOWLEDGE_COLLECTION}),
        encoding="utf-8"
    )

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s %(levelname)s %(message)s",
    handlers=[
        logging.FileHandler(r"D:\OneDrive\Claude\agents\watson\logs\kb_watcher.log"),
        logging.StreamHandler()
    ]
)
log = logging.getLogger(__name__)


def convert_docx(path: Path) -> str:
    import docx
    doc = docx.Document(str(path))
    lines = []
    for para in doc.paragraphs:
        if para.style.name.startswith("Heading"):
            level = para.style.name[-1] if para.style.name[-1].isdigit() else "1"
            lines.append(f"{'#' * int(level)} {para.text}")
        elif para.text.strip():
            lines.append(para.text)
    return "\n\n".join(lines)


def convert_pdf(path: Path) -> str:
    import pdfplumber
    pages = []
    with pdfplumber.open(str(path)) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
    return "\n\n".join(pages)


def convert_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        lines = [f"## Slide {i}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
        slides.append("\n".join(lines))
    return "\n\n".join(slides)


def convert_xlsx(path: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(str(path), data_only=True)
    sheets = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        lines = [f"## Sheet: {sheet}"]
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                lines.append(" | ".join(cells))
        sheets.append("\n".join(lines))
    return "\n\n".join(sheets)


def convert_csv(path: Path) -> str:
    import csv
    rows = []
    with open(path, encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            rows.append(" | ".join(row))
    return "\n".join(rows)


def process_file(path: Path):
    suffix = path.suffix.lower()
    stem = path.stem

    try:
        if suffix in (".txt", ".md"):
            content = path.read_text(encoding="utf-8", errors="ignore")
            out_suffix = suffix
        elif suffix == ".docx":
            content = convert_docx(path)
            out_suffix = ".md"
        elif suffix == ".pdf":
            content = convert_pdf(path)
            out_suffix = ".txt"
        elif suffix == ".pptx":
            content = convert_pptx(path)
            out_suffix = ".md"
        elif suffix in (".xlsx", ".xls"):
            content = convert_xlsx(path)
            out_suffix = ".txt"
        elif suffix == ".csv":
            content = convert_csv(path)
            out_suffix = ".txt"
        else:
            log.warning("Unsupported file type: %s — skipping", path.name)
            return

        # Save to library
        out_path = LIBRARY / (stem + out_suffix)
        # Handle duplicates
        counter = 1
        while out_path.exists():
            out_path = LIBRARY / (f"{stem}_{counter}{out_suffix}")
            counter += 1

        out_path.write_text(content, encoding="utf-8")
        log.info("Converted: %s -> %s", path.name, out_path.name)

        # Ingest into OpenWebUI
        ingest_file(out_path)

        # Delete original from inbox
        path.unlink()
        log.info("Deleted original: %s", path.name)

    except Exception as e:
        log.error("Failed to process %s: %s", path.name, e)


def ingest_file(path: Path):
    """Upload file to OpenWebUI knowledge base."""
    if not OPENWEBUI_API_KEY:
        log.warning("No OPENWEBUI_API_KEY — skipping ingest")
        return
    try:
        # Upload file
        with open(path, "rb") as f:
            resp = requests.post(
                f"{OPENWEBUI_URL}/api/v1/files/",
                headers={"Authorization": f"Bearer {OPENWEBUI_API_KEY}"},
                files={"file": (path.name, f, "text/plain")},
                timeout=30
            )
        if resp.status_code != 200:
            log.error("Upload failed: %s %s", resp.status_code, resp.text)
            return
        file_id = resp.json().get("id")
        log.info("Uploaded to OpenWebUI: %s (id=%s)", path.name, file_id)

        global _collection_id

        if not _collection_id:
            # 1. Check local cache — survives restarts and works around unreliable GET /knowledge/
            _collection_id = _load_cached_collection_id()
            if _collection_id:
                log.info("Loaded collection id from cache: %s", _collection_id)

        if not _collection_id:
            # 2. Fall back to querying the API (handles both list and paginated {items, total} shapes)
            cols_resp = requests.get(
                f"{OPENWEBUI_URL}/api/v1/knowledge/",
                headers={"Authorization": f"Bearer {OPENWEBUI_API_KEY}"},
                timeout=10
            )
            if cols_resp.status_code == 200:
                data = cols_resp.json()
                cols = data.get("items", data) if isinstance(data, dict) else data
            else:
                cols = []

            if isinstance(cols, list):
                _collection_id = next((c["id"] for c in cols if isinstance(c, dict) and c.get("name") == KNOWLEDGE_COLLECTION), None)
            if _collection_id:
                log.info("Found existing collection via API: %s", _collection_id)
                _save_cached_collection_id(_collection_id)

        if not _collection_id:
            # 3. Create a new collection and cache its ID immediately
            col_resp = requests.post(
                f"{OPENWEBUI_URL}/api/v1/knowledge/create",
                headers={
                    "Authorization": f"Bearer {OPENWEBUI_API_KEY}",
                    "Content-Type": "application/json"
                },
                json={"name": KNOWLEDGE_COLLECTION, "description": "Personal knowledge library"},
                timeout=10
            )
            if col_resp.status_code == 200:
                _collection_id = col_resp.json().get("id")
                log.info("Created collection '%s' (id=%s)", KNOWLEDGE_COLLECTION, _collection_id)
                _save_cached_collection_id(_collection_id)

        if _collection_id and file_id:
            add_resp = requests.post(
                f"{OPENWEBUI_URL}/api/v1/knowledge/{_collection_id}/file/add",
                headers={
                    "Authorization": f"Bearer {OPENWEBUI_API_KEY}",
                    "Content-Type": "application/json"
                },
                json={"file_id": file_id},
                timeout=10
            )
            if add_resp.status_code == 200:
                log.info("Added to collection '%s'", KNOWLEDGE_COLLECTION)
            else:
                log.error("Failed to add to collection: %s %s", add_resp.status_code, add_resp.text)
        else:
            log.warning("Could not find or create collection '%s'", KNOWLEDGE_COLLECTION)

    except Exception as e:
        log.error("Ingest failed for %s: %s", path.name, e)


def watch():
    log.info("KB Watcher started. Monitoring: %s", INBOX)
    INBOX.mkdir(parents=True, exist_ok=True)
    LIBRARY.mkdir(parents=True, exist_ok=True)
    seen = set()

    while True:
        try:
            for path in INBOX.iterdir():
                if path.is_file() and path not in seen:
                    # Wait briefly to ensure file is fully written
                    time.sleep(2)
                    if path.exists():
                        seen.add(path)
                        process_file(path)
            # Clean seen set of files no longer in inbox
            seen = {p for p in seen if p.exists()}
        except Exception as e:
            log.error("Watcher error: %s", e)
        time.sleep(5)


if __name__ == "__main__":
    watch()
