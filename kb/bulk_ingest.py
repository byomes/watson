"""
bulk_ingest.py — Re-populate the Personal Library knowledge collection in OpenWebUI.

Targets two directories:
  1. LIBRARY (F:\\Knowledge_Database\\_library) — already-converted txt/md files.
     These are uploaded directly without conversion.
  2. INBOX   (F:\\Knowledge_Database\\_inbox)   — raw incoming files.
     These are converted (same logic as watcher.py) then uploaded.

Usage:
    python kb/bulk_ingest.py [--library-only] [--inbox-only]
"""

import argparse
import csv
import json
import os
import sys
from pathlib import Path

import requests
from dotenv import load_dotenv

load_dotenv(r"D:\OneDrive\Claude\agents\watson\.env")

INBOX = Path(r"F:\\Knowledge_Database\\_pending_ingest")
LIBRARY = Path(r"F:\Knowledge_Database\_library")
OPENWEBUI_URL = os.getenv("OPENWEBUI_URL", "http://localhost:3000")
OPENWEBUI_API_KEY = os.getenv("OPENWEBUI_API_KEY")
CACHE_FILE = Path(__file__).parent / ".collection_id_cache.json"
DEFAULT_COLLECTION_ID = "62aa6542-0513-4d60-900d-85f386a64f13"

SUPPORTED_INBOX_TYPES = {".txt", ".md", ".docx", ".pdf", ".pptx", ".xlsx", ".xls", ".csv"}


def load_collection_id() -> str:
    try:
        if CACHE_FILE.exists():
            data = json.loads(CACHE_FILE.read_text(encoding="utf-8"))
            cid = data.get("collection_id")
            if cid:
                return cid
    except Exception:
        pass
    return DEFAULT_COLLECTION_ID


# ── converters (mirrors watcher.py) ──────────────────────────────────────────

def convert_docx(path: Path) -> str:
    import docx
    doc = docx.Document(str(path))
    lines = []
    for para in doc.paragraphs:
        if para.style.name.startswith("Heading"):
            level = para.style.name[-1] if para.style.name[-1].isdigit() else "1"
            lines.append(f"{'#' * int(level)} {para.text}")
        elif para.text.strip():
            lines.append(para.text)
    return "\n\n".join(lines)


def convert_pdf(path: Path) -> str:
    import pdfplumber
    pages = []
    with pdfplumber.open(str(path)) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
    return "\n\n".join(pages)


def convert_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        lines = [f"## Slide {i}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
        slides.append("\n".join(lines))
    return "\n\n".join(slides)


def convert_xlsx(path: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(str(path), data_only=True)
    sheets = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        lines = [f"## Sheet: {sheet}"]
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                lines.append(" | ".join(cells))
        sheets.append("\n".join(lines))
    return "\n\n".join(sheets)


def convert_csv(path: Path) -> str:
    rows = []
    with open(path, encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            rows.append(" | ".join(row))
    return "\n".join(rows)


def convert_inbox_file(path: Path) -> tuple[str, str]:
    """Return (content, out_suffix). Raises on unsupported type."""
    suffix = path.suffix.lower()
    if suffix in (".txt", ".md"):
        return path.read_text(encoding="utf-8", errors="ignore"), suffix
    if suffix == ".docx":
        return convert_docx(path), ".md"
    if suffix == ".pdf":
        return convert_pdf(path), ".txt"
    if suffix == ".pptx":
        return convert_pptx(path), ".md"
    if suffix in (".xlsx", ".xls"):
        return convert_xlsx(path), ".txt"
    if suffix == ".csv":
        return convert_csv(path), ".txt"
    raise ValueError(f"Unsupported file type: {suffix}")


# ── OpenWebUI helpers ─────────────────────────────────────────────────────────

def upload_file(path: Path, headers: dict) -> str | None:
    """Upload a text/md file to OpenWebUI and return its file_id."""
    with open(path, "rb") as f:
        resp = requests.post(
            f"{OPENWEBUI_URL}/api/v1/files/",
            headers=headers,
            files={"file": (path.name, f, "text/plain")},
            timeout=60,
        )
    if resp.status_code != 200:
        print(f"    ERROR upload failed: {resp.status_code} {resp.text[:200]}")
        return None
    return resp.json().get("id")


def add_to_collection(file_id: str, collection_id: str, headers: dict) -> bool:
    resp = requests.post(
        f"{OPENWEBUI_URL}/api/v1/knowledge/{collection_id}/file/add",
        headers={**headers, "Content-Type": "application/json"},
        json={"file_id": file_id},
        timeout=30,
    )
    if resp.status_code == 200:
        return True
    print(f"    ERROR add-to-collection: {resp.status_code} {resp.text[:200]}")
    return False


# ── ingest runners ────────────────────────────────────────────────────────────

def ingest_library(collection_id: str, headers: dict) -> tuple[int, int]:
    """Upload every txt/md file in LIBRARY directly. Returns (ok, failed)."""
    files = sorted(p for p in LIBRARY.iterdir() if p.is_file() and p.suffix.lower() in (".txt", ".md"))
    if not files:
        print("LIBRARY: no files found.")
        return 0, 0

    print(f"\nLIBRARY — {len(files)} file(s) to ingest from {LIBRARY}\n")
    ok = failed = 0
    for i, path in enumerate(files, 1):
        print(f"  [{i}/{len(files)}] {path.name}")
        file_id = upload_file(path, headers)
        if not file_id:
            failed += 1
            continue
        print(f"    uploaded  file_id={file_id}")
        if add_to_collection(file_id, collection_id, headers):
            print(f"    added to collection")
            ok += 1
        else:
            failed += 1

    return ok, failed


def ingest_inbox(collection_id: str, headers: dict) -> tuple[int, int]:
    """Convert + upload every supported file in INBOX. Returns (ok, failed)."""
    files = sorted(
        p for p in INBOX.iterdir()
        if p.is_file() and p.suffix.lower() in SUPPORTED_INBOX_TYPES
    )
    if not files:
        print("INBOX: no files found.")
        return 0, 0

    print(f"\nINBOX — {len(files)} file(s) to convert and ingest from {INBOX}\n")
    ok = failed = 0
    for i, path in enumerate(files, 1):
        print(f"  [{i}/{len(files)}] {path.name}")
        try:
            content, out_suffix = convert_inbox_file(path)
        except Exception as e:
            print(f"    ERROR conversion failed: {e}")
            failed += 1
            continue

        # Save converted file to LIBRARY
        out_path = LIBRARY / (path.stem + out_suffix)
        counter = 1
        while out_path.exists():
            out_path = LIBRARY / (f"{path.stem}_{counter}{out_suffix}")
            counter += 1
        out_path.write_text(content, encoding="utf-8")
        print(f"    converted -> {out_path.name}")

        file_id = upload_file(out_path, headers)
        if not file_id:
            failed += 1
            continue
        print(f"    uploaded  file_id={file_id}")
        if add_to_collection(file_id, collection_id, headers):
            print(f"    added to collection")
            ok += 1
        else:
            failed += 1

    return ok, failed


# ── main ──────────────────────────────────────────────────────────────────────

def main():
    if not OPENWEBUI_API_KEY:
        sys.exit("OPENWEBUI_API_KEY not set — aborting")

    parser = argparse.ArgumentParser(description="Bulk ingest files into OpenWebUI Personal Library")
    parser.add_argument("--library-only", action="store_true", help="Only ingest LIBRARY files")
    parser.add_argument("--inbox-only", action="store_true", help="Only process INBOX files")
    args = parser.parse_args()

    collection_id = load_collection_id()
    headers = {"Authorization": f"Bearer {OPENWEBUI_API_KEY}"}

    print(f"Collection ID : {collection_id}")
    print(f"OpenWebUI URL : {OPENWEBUI_URL}")

    total_ok = total_failed = 0

    if not args.inbox_only:
        ok, failed = ingest_library(collection_id, headers)
        total_ok += ok
        total_failed += failed

    if not args.library_only:
        ok, failed = ingest_inbox(collection_id, headers)
        total_ok += ok
        total_failed += failed

    print(f"\n{'─' * 50}")
    print(f"Done. {total_ok} ingested, {total_failed} failed.")


if __name__ == "__main__":
    main()
